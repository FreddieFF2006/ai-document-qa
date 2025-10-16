import streamlit as st
from pdf_ai_agent import EmbeddingManager, ClaudeAIAgent
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
import re
import os
from dotenv import load_dotenv
import uuid
from datetime import datetime

# Load environment variables from .env file
load_dotenv()

st.set_page_config(page_title="AI Document Q&A", page_icon="🤖", layout="wide")

# Add custom CSS to hide the ugly file uploader and create clean UI
st.markdown("""
<style>
    /* Center the icons in buttons */
    .stButton button {
        display: flex;
        align-items: center;
        justify-content: center;
        padding: 0.25rem 0.75rem;
    }
    
    /* Make chat buttons take full width */
    div[data-testid="column"] button {
        width: 100%;
    }
    
    /* Force trash can button to be perfectly centered */
    div[data-testid="column"]:last-child button {
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        padding: 0.25rem !important;
        min-width: 40px !important;
    }
    
    /* Ensure emoji stays centered */
    div[data-testid="column"]:last-child button p {
        display: flex;
        align-items: center;
        justify-content: center;
        margin: 0 !important;
        padding: 0 !important;
        line-height: 1;
    }
    
    /* HIDE THE ENTIRE FILE UPLOADER UI */
    section[data-testid="stFileUploader"] {
        display: none !important;
    }
    
    /* Style for the custom attach button */
    .attach-button {
        background-color: #f0f2f6;
        border: 1px solid #e0e0e0;
        border-radius: 8px;
        padding: 8px 12px;
        cursor: pointer;
        display: inline-flex;
        align-items: center;
        justify-content: center;
        font-size: 20px;
        transition: all 0.2s;
        width: 40px;
        height: 40px;
    }
    
    .attach-button:hover {
        background-color: #e0e0e0;
        border-color: #d0d0d0;
    }
    
    /* Fix sidebar button alignment */
    .stSidebar button {
        display: flex;
        align-items: center;
        justify-content: center;
    }
</style>
""", unsafe_allow_html=True)

def smart_chunk_text(text, chunk_size=1500, overlap=100):
    """
    Optimized chunking with smaller chunks for maximum precision
    
    Args:
        chunk_size: Target size for each chunk (default: 1500)
        overlap: Minimal overlap for context (default: 100)
    
    Features:
    - Smaller chunks = more precise retrieval
    - Better granularity for answering specific questions
    - E5-Large-v2 handles these well
    - Preserves paragraph boundaries
    """
    
    # Split by paragraphs
    paragraphs = re.split(r'\n\s*\n+', text)
    
    chunks = []
    current_chunk = ""
    
    for para in paragraphs:
        para = para.strip()
        if not para or len(para) < 50:  # Skip very short paragraphs
            continue
        
        # If adding this paragraph is within size limit, add it
        if len(current_chunk) + len(para) + 2 <= chunk_size:
            if current_chunk:
                current_chunk += "\n\n" + para
            else:
                current_chunk = para
        else:
            # Save current chunk if substantial
            if len(current_chunk) > 300:  # Lowered threshold for smaller chunks
                chunks.append(current_chunk)
            
            # Handle oversized paragraphs
            if len(para) > chunk_size:
                # Split by sentences only for huge paragraphs
                sentences = re.split(r'(?<=[.!?])\s+', para)
                temp_chunk = ""
                
                for sentence in sentences:
                    if len(temp_chunk) + len(sentence) + 1 <= chunk_size:
                        temp_chunk += (" " if temp_chunk else "") + sentence
                    else:
                        if len(temp_chunk) > 300:
                            chunks.append(temp_chunk)
                        temp_chunk = sentence
                
                current_chunk = temp_chunk
            else:
                current_chunk = para
    
    # Add final chunk if substantial
    if len(current_chunk.strip()) > 300:
        chunks.append(current_chunk.strip())
    
    # Aggressive duplicate removal
    unique_chunks = []
    seen = set()
    for chunk in chunks:
        # Use first 100 chars as fingerprint
        fingerprint = chunk[:100].lower().strip()
        if fingerprint not in seen and len(chunk) > 300:
            seen.add(fingerprint)
            unique_chunks.append(chunk)
    
    print(f"✅ Created {len(unique_chunks)} chunks from {len(text):,} characters")
    
    return unique_chunks

def extract_text_from_pdf(file):
    """Extract text from PDF with multiple fallback methods"""
    text = ""
    
    try:
        # Method 1: Try PyPDF2 first (fastest)
        reader = PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        
        # If we got substantial text, return it
        if len(text.strip()) > 100:
            return text
    except Exception as e:
        print(f"PyPDF2 failed: {e}")
    
    # Method 2: Try pdfplumber (more robust)
    try:
        import pdfplumber
        file.seek(0)  # Reset file pointer
        
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        if len(text.strip()) > 100:
            return text
    except ImportError:
        print("pdfplumber not installed")
    except Exception as e:
        print(f"pdfplumber failed: {e}")
    
    # Method 3: Try pymupdf/fitz (most powerful)
    try:
        import fitz  # PyMuPDF
        file.seek(0)  # Reset file pointer
        
        # Read file bytes
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        for page in doc:
            page_text = page.get_text()
            if page_text:
                text += page_text + "\n"
        
        doc.close()
        
        if len(text.strip()) > 100:
            return text
    except ImportError:
        print("PyMuPDF not installed")
    except Exception as e:
        print(f"PyMuPDF failed: {e}")
    
    # If all methods failed
    if len(text.strip()) < 100:
        raise Exception("Could not extract text from PDF. File may be corrupted, image-based, or encrypted.")
    
    return text

def extract_text_from_docx(file):
    """Extract text from Word document"""
    doc = DocxDocument(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pptx(file):
    """Extract text from PowerPoint"""
    try:
        from pptx import Presentation
        prs = Presentation(file)
        all_text = ""
        for i, slide in enumerate(prs.slides, 1):
            all_text += f"\n=== Slide {i} ===\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text += shape.text + "\n"
            all_text += "\n"
        return all_text
    except:
        return ""

def extract_text_from_excel(file):
    """Extract text from Excel"""
    try:
        import pandas as pd
        excel_file = pd.ExcelFile(file)
        all_text = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file, sheet_name=sheet_name)
            all_text += f"\n=== Sheet: {sheet_name} ===\n\n"
            all_text += df.to_string(index=False) + "\n\n"
        return all_text
    except ImportError:
        return ""

def extract_text_from_txt(file):
    """Extract text from text file"""
    return file.read().decode('utf-8')

def process_uploaded_files(uploaded_files, chat_id):
    """Process uploaded files and return chunks with progress tracking"""
    current_chat = st.session_state.chats[chat_id]
    
    if not current_chat['embedding_manager']:
        current_chat['embedding_manager'] = EmbeddingManager()
    
    all_chunks = []
    processed_files = []
    total_chars = 0
    
    # Progress tracking
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    total_files = len(uploaded_files)
    
    for idx, uploaded_file in enumerate(uploaded_files):
        file_ext = uploaded_file.name.split('.')[-1].lower()
        
        # Update progress
        progress = (idx + 1) / total_files
        progress_bar.progress(progress)
        status_text.text(f"Processing {uploaded_file.name} ({idx + 1}/{total_files})...")
        
        try:
            if file_ext == 'pdf':
                text = extract_text_from_pdf(uploaded_file)
            elif file_ext in ['docx', 'doc']:
                text = extract_text_from_docx(uploaded_file)
            elif file_ext in ['pptx', 'ppt']:
                text = extract_text_from_pptx(uploaded_file)
            elif file_ext == 'txt':
                text = extract_text_from_txt(uploaded_file)
            else:
                continue
            
            if text and len(text.strip()) > 0:
                total_chars += len(text)
                
                # Show chunking progress
                status_text.text(f"Creating precise chunks for {uploaded_file.name}...")
                chunks = smart_chunk_text(text, chunk_size=1500, overlap=100)
                
                chunks_with_source = [f"[Source: {uploaded_file.name}]\n\n{chunk}" for chunk in chunks]
                all_chunks.extend(chunks_with_source)
                processed_files.append(uploaded_file.name)
        
        except Exception as e:
            st.warning(f"⚠️ Could not process {uploaded_file.name}: {str(e)}")
    
    # Clear progress indicators
    progress_bar.empty()
    status_text.empty()
    
    if all_chunks:
        # Show embedding progress
        with st.spinner(f"Creating E5-Large embeddings for {len(all_chunks)} chunks..."):
            current_chat['chunks'].extend(all_chunks)
            current_chat['embedding_manager'].build_index(current_chat['chunks'])
            current_chat['processed_files'].extend(processed_files)
        
        return len(processed_files), len(all_chunks), total_chars
    
    return 0, 0, 0

def generate_chat_title(first_message):
    """Generate a short title from the first message"""
    title = first_message[:40]
    if len(first_message) > 40:
        title += "..."
    return title

def create_new_chat():
    """Create a new chat session"""
    chat_id = str(uuid.uuid4())
    st.session_state.chats[chat_id] = {
        'title': 'New Chat',
        'messages': [],
        'chunks': [],
        'embedding_manager': None,
        'processed_files': [],
        'created_at': datetime.now()
    }
    st.session_state.current_chat_id = chat_id
    return chat_id

# Check for API key
api_key = os.getenv('ANTHROPIC_API_KEY')

if not api_key:
    st.error("❌ API Key not found!")
    st.info("Please make sure your .env file contains: ANTHROPIC_API_KEY=your-key-here")
    st.stop()

# Initialize Claude agent (shared across all chats)
if 'claude_agent' not in st.session_state:
    try:
        st.session_state.claude_agent = ClaudeAIAgent(api_key)
    except Exception as e:
        st.error(f"Error initializing Claude: {e}")
        st.stop()

# Initialize chats dictionary
if 'chats' not in st.session_state:
    st.session_state.chats = {}
    create_new_chat()

# Initialize current chat ID
if 'current_chat_id' not in st.session_state:
    if st.session_state.chats:
        st.session_state.current_chat_id = list(st.session_state.chats.keys())[0]
    else:
        create_new_chat()

# Initialize show uploader state
if 'show_file_uploader' not in st.session_state:
    st.session_state.show_file_uploader = False

# Header
st.title("🤖 AI Document Q&A Assistant")
st.markdown("Upload your documents and ask questions! **Powered by Claude Opus 4** 🧠")

# Sidebar - Chat management
with st.sidebar:
    st.header("💬 Chats")
    
    # New chat button
    if st.button("➕ New Chat", use_container_width=True):
        create_new_chat()
        st.rerun()
    
    st.markdown("---")
    
    # List all chats
    chat_items = sorted(st.session_state.chats.items(), 
                       key=lambda x: x[1]['created_at'], 
                       reverse=True)
    
    for chat_id, chat_data in chat_items:
        col1, col2 = st.columns([5, 1])
        
        with col1:
            # Chat button
            is_current = chat_id == st.session_state.current_chat_id
            button_type = "primary" if is_current else "secondary"
            icon = "📍" if is_current else "💬"
            
            if st.button(
                f"{icon} {chat_data['title']}", 
                key=f"chat_{chat_id}",
                use_container_width=True,
                type=button_type
            ):
                st.session_state.current_chat_id = chat_id
                st.rerun()
        
        with col2:
            # Delete button - centered
            if st.button("🗑️", key=f"delete_{chat_id}", use_container_width=True):
                if len(st.session_state.chats) > 1:
                    del st.session_state.chats[chat_id]
                    # Switch to another chat
                    if st.session_state.current_chat_id == chat_id:
                        st.session_state.current_chat_id = list(st.session_state.chats.keys())[0]
                    st.rerun()
                else:
                    st.warning("Cannot delete the last chat!")

# Get current chat
current_chat = st.session_state.chats[st.session_state.current_chat_id]

# Main chat interface
st.header(f"💬 {current_chat['title']}")

# Show document count for current chat
if current_chat['processed_files']:
    st.caption(f"📚 {len(current_chat['processed_files'])} document(s) loaded | {len(current_chat['chunks'])} chunks | 🧠 Claude Opus 4")

# Display chat history
for message in current_chat['messages']:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Chat input area - clean design
input_container = st.container()

with input_container:
    col1, col2 = st.columns([0.6, 9.4])
    
    with col1:
        # Clean attach button that triggers file uploader
        if st.button("➕", key="attach_btn", help="Attach files"):
            st.session_state.show_file_uploader = not st.session_state.show_file_uploader
    
    with col2:
        prompt = st.chat_input("Ask me anything...")

# Show file uploader in a clean way when button is clicked
if st.session_state.show_file_uploader:
    uploaded_files = st.file_uploader(
        "Choose files to upload",
        type=['pdf', 'docx', 'doc', 'txt', 'pptx', 'ppt'],
        accept_multiple_files=True,
        key=f"uploader_{st.session_state.current_chat_id}"
    )
    
    if uploaded_files:
        new_files = [f for f in uploaded_files if f.name not in current_chat['processed_files']]
        
        if new_files:
            with st.spinner(f"Processing {len(new_files)} file(s)..."):
                num_files, num_chunks, total_chars = process_uploaded_files(new_files, st.session_state.current_chat_id)
                
                if num_files > 0:
                    st.success(f"✅ Processed {num_files} file(s): {num_chunks} chunks, {total_chars:,} characters")
                    st.session_state.show_file_uploader = False
                    st.rerun()
                else:
                    st.error("Could not extract text from the uploaded files")

# Handle chat input - MAXIMUM CHUNK RETRIEVAL WITH OPUS
if prompt:
    # Update chat title if this is the first message
    if len(current_chat['messages']) == 0:
        current_chat['title'] = generate_chat_title(prompt)
    
    # Add user message to chat
    current_chat['messages'].append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)
    
    # Generate response
    with st.chat_message("assistant"):
        with st.spinner("🧠 Claude Opus 4 thinking..."):
            try:
                # Check if we have documents
                if current_chat['embedding_manager'] and current_chat['chunks']:
                    # RETRIEVE 25 CHUNKS (increased from 15 for maximum coverage)
                    results = current_chat['embedding_manager'].search(prompt, top_k=25)
                    chunks = [chunk for chunk, score in results]
                    
                    # Calculate approximate tokens
                    total_chars = sum(len(c) for c in chunks)
                    approx_tokens = total_chars // 4
                    
                    # Show search info with more details
                    with st.expander(f"📊 Search: {len(chunks)} chunks, ~{approx_tokens:,} tokens"):
                        st.caption(f"Retrieved {len(chunks)} most relevant chunks")
                        st.caption(f"Context size: {total_chars:,} characters ≈ {approx_tokens:,} tokens")
                        st.caption("🔍 Maximum coverage mode - searching through extensive context")
                        st.caption("🧠 **Powered by Claude Opus 4** - Maximum Intelligence")
                        
                        # Show top 5 chunk scores
                        st.caption("**Top 5 relevance scores:**")
                        for i, (chunk, score) in enumerate(results[:5], 1):
                            # Extract first 60 chars as preview
                            preview = chunk.replace("[Source:", "").strip()[:60] + "..."
                            st.caption(f"{i}. Score: {score:.3f} - {preview}")
                    
                    # Get answer from Claude Opus WITH CONVERSATION MEMORY
                    answer = st.session_state.claude_agent.ask(
                        prompt, 
                        chunks, 
                        max_tokens=4000,  # Increased from 3000 for longer answers
                        conversation_history=current_chat['messages']
                    )
                    
                else:
                    # No documents - just chat with Claude Opus directly
                    message = st.session_state.claude_agent.client.messages.create(
                        model="claude-opus-4-20250514",  # ✅ CLAUDE OPUS 4
                        max_tokens=4000,
                        messages=[{"role": "user", "content": prompt}]
                    )
                    answer = message.content[0].text
                
                st.markdown(answer)
                current_chat['messages'].append({"role": "assistant", "content": answer})
                
            except Exception as e:
                error_msg = f"Error: {str(e)}"
                st.error(error_msg)
                current_chat['messages'].append({"role": "assistant", "content": error_msg})